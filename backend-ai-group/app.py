# ============================================
# app.py — SIMULATION ENGINE MODULE (FIXED VISION & UPLOAD)
# ============================================

import json
import httpx
import os
import asyncio
import re
import math
import io
import base64
from fastapi import FastAPI, WebSocket, WebSocketDisconnect, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from dotenv import load_dotenv
from openpyxl import load_workbook
load_dotenv()

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

os.makedirs("static", exist_ok=True)
app.mount("/static", StaticFiles(directory="static"), name="static")

api_key = os.getenv("GROQ_API_KEY")
TARGET_MODEL = "llama-3.1-8b-instant"
VISION_MODEL = "meta-llama/llama-4-scout-17b-16e-instruct"
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB
MAX_EXTRACTED_TEXT_LENGTH = 8000    # Maks karakter teks yang diekstrak

THEME = {
    "primary": RGBColor(70, 72, 212),
    "secondary": RGBColor(96, 99, 238),
    "accent": RGBColor(16, 185, 129),
    "dark": RGBColor(30, 41, 59),
    "text": RGBColor(51, 65, 85),
    "light": RGBColor(248, 250, 252),
    "white": RGBColor(255, 255, 255),
}


# ==================== DOCUMENT READER ====================

def extract_from_docx(content: bytes) -> str:
    try:
        from docx import Document
    except ImportError:
        return "[ERROR] Library python-docx belum terinstall. Jalankan: pip install python-docx"

    try:
        doc = Document(io.BytesIO(content))
        text_parts = []

        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text.strip())

        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    text_parts.append(" | ".join(row_text))

        for section in doc.sections:
            header = section.header
            if header and header.paragraphs:
                h_text = " ".join(p.text.strip() for p in header.paragraphs if p.text.strip())
                if h_text:
                    text_parts.append(f"[Header: {h_text}]")
            footer = section.footer
            if footer and footer.paragraphs:
                f_text = " ".join(p.text.strip() for p in footer.paragraphs if p.text.strip())
                if f_text:
                    text_parts.append(f"[Footer: {f_text}]")

        result = "\n".join(text_parts)
        if not result.strip():
            return "[Tidak ada teks ditemukan dalam dokumen Word]"
        return result
    except Exception as e:
        error_msg = str(e).lower()
        if "not a zip file" in error_msg or "bad magic" in error_msg:
            return "[ERROR] Format .doc lama tidak didukung. Silakan simpan ulang sebagai .docx (Word 2007+)"
        return f"[ERROR membaca dokumen Word: {str(e)}]"


def extract_from_xlsx(content: bytes) -> str:
    try:
        import openpyxl
    except ImportError:
        return "[ERROR] Library openpyxl belum terinstall. Jalankan: pip install openpyxl"

    try:
        wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
        text_parts = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"=== Sheet: {sheet_name} ===")

            if ws.max_row is None or ws.max_column is None:
                text_parts.append("[Sheet kosong]")
                continue

            headers = []
            data_rows = []

            for row_idx, row in enumerate(ws.iter_rows(values_only=True), 1):
                row_text = [str(cell).strip() if cell is not None else "" for cell in row]
                if not any(t for t in row_text):
                    continue

                if row_idx == 1:
                    headers = [t for t in row_text if t]
                    if headers:
                        text_parts.append("Kolom: " + " | ".join(headers))
                    else:
                        text_parts.append(" | ".join(row_text))
                else:
                    if headers and len(row_text) >= len(headers):
                        formatted = []
                        for h, v in zip(headers, row_text):
                            if v:
                                formatted.append(f"{h}: {v}")
                        if formatted:
                            data_rows.append("; ".join(formatted))
                    else:
                        data_rows.append(" | ".join(t for t in row_text if t))

            if data_rows:
                text_parts.extend(data_rows)

        result = "\n".join(text_parts)
        if not result.strip() or result.strip() == "=" * 20:
            return "[Tidak ada data ditemukan dalam file Excel]"
        return result
    except Exception as e:
        error_msg = str(e).lower()
        if "not a zip file" in error_msg:
            return "[ERROR] Format .xls lama tidak didukung. Silakan simpan ulang sebagai .xlsx"
        return f"[ERROR membaca file Excel: {str(e)}]"


def extract_from_pdf(content: bytes) -> str:
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        try:
            from pdfplumber import PDF as PdfReader_plumber
            return _extract_pdf_with_pdfplumber(content)
        except ImportError:
            return "[ERROR] Library PDF belum terinstall. Jalankan: pip install PyPDF2"

    try:
        reader = PdfReader(io.BytesIO(content))
        text_parts = []

        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():
                text_parts.append(f"--- Halaman {i+1} ---")
                text_parts.append(text.strip())

        result = "\n".join(text_parts)
        if not result.strip():
            return "[Tidak ada teks ditemukan dalam PDF. Mungkin PDF berisi gambar scanned — coba upload sebagai gambar (PNG/JPG)]"
        return result
    except Exception as e:
        return f"[ERROR membaca PDF: {str(e)}]"


def _extract_pdf_with_pdfplumber(content: bytes) -> str:
    try:
        import pdfplumber
    except ImportError:
        return "[pdfplumber tidak tersedia]"

    text_parts = []
    with pdfplumber.open(io.BytesIO(content)) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text and text.strip():
                text_parts.append(f"--- Halaman {i+1} ---")
                text_parts.append(text.strip())

            tables = page.extract_tables()
            for table in tables:
                for row in table:
                    row_text = [str(cell).strip() if cell else "" for cell in row]
                    if any(row_text):
                        text_parts.append(" | ".join(row_text))

    result = "\n".join(text_parts)
    if not result.strip():
        return "[Tidak ada teks ditemukan dalam PDF]"
    return result


async def extract_from_image(content: bytes, filename: str) -> str:
    """
    Ekstrak teks dari gambar (PNG/JPG) menggunakan Groq Vision API.
    Diperbaiki: Resize agresif ke 1024px, kompresi JPEG, dan error handling spesifik.
    """
    try:
        from PIL import Image
    except ImportError:
        return "[ERROR] Library Pillow belum terinstall. Jalankan: pip install Pillow"

    try:
        # 1. Buka gambar dan konversi ke RGB (hapus alpha channel dari PNG)
        img = Image.open(io.BytesIO(content))
        if img.mode in ('RGBA', 'P', 'LA'):
            img = img.convert('RGB')

        # 2. Resize gambar (Max 1024px) untuk menghindari payload terlalu besar yang menyebabkan Error 400
        max_dim = 1024
        if max(img.size) > max_dim:
            ratio = max_dim / max(img.size)
            new_size = (int(img.size[0] * ratio), int(img.size[1] * ratio))
            img = img.resize(new_size, Image.Resampling.LANCZOS)

        # 3. Simpan ke buffer dalam format JPEG dengan kualitas 80 untuk kompresi ukuran base64
        buf = io.BytesIO()
        img.save(buf, format='JPEG', quality=80)
        b64_image = base64.b64encode(buf.getvalue()).decode('utf-8')

        # 4. Siapkan payload untuk Groq Vision API
        headers = {
            "Authorization": f"Bearer {GROQ_API_KEY}",
            "Content-Type": "application/json"
        }

        messages = [{
            "role": "user",
            "content": [
                {
                    "type": "text",
                    "text": "Baca dan ekstrak SEMUA teks, angka, tabel, label, dan data dari gambar ini. Sertakan setiap nilai numerik, keterangan, dan deskripsi. Jika ada tabel, format ulang sebagai teks terstruktur. Jika ini dokumen/surat, baca seluruh isinya secara lengkap dan akurat. Output dalam teks biasa yang terstruktur."
                },
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/jpeg;base64,{b64_image}"
                    }
                }
            ]
        }]

        # 5. Kirim request ke Groq
        async with httpx.AsyncClient(timeout=60.0) as client:
            res = await client.post(
                "https://api.groq.com/openai/v1/chat/completions",
                json={
                    "model": VISION_MODEL, 
                    "messages": messages, 
                    "max_tokens": 1500, 
                    "temperature": 0.1
                },
                headers=headers,
            )
            
            if res.status_code == 200:
                res_json = res.json()
                if "choices" in res_json and len(res_json["choices"]) > 0:
                    extracted = res_json["choices"][0]["message"]["content"]
                    if extracted and extracted.strip():
                        return extracted.strip()
                    return "[Tidak ada teks terdeteksi dalam gambar]"
            elif res.status_code == 429:
                return "[Rate limit tercapai saat OCR gambar. Coba lagi dalam beberapa detik.]"
            else:
                # Tangkap error sebenarnya dari Groq untuk debugging
                error_detail = res.text[:500]
                print(f"💥 [VISION API ERROR] Status: {res.status_code} | Detail: {error_detail}")
                return f"[Gagal membaca gambar. Status: {res.status_code}. Detail: {error_detail}]"

    except Exception as e:
        return f"[ERROR memproses gambar: {str(e)}]"


def truncate_extracted_text(text: str, max_length: int = MAX_EXTRACTED_TEXT_LENGTH) -> str:
    if len(text) <= max_length:
        return text
    return text[:max_length] + f"\n\n[... DOKUMEN DIPOTONG — Total {len(text)} karakter, hanya {max_length} pertama yang ditampilkan ...]"


@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    if not file.filename:
        raise HTTPException(status_code=400, detail="Tidak ada nama file")

    content = await file.read()

    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(status_code=413, detail=f"File terlalu besar. Ukuran maksimal: {MAX_FILE_SIZE // (1024*1024)}MB")

    filename = file.filename.lower()
    extracted_text = ""
    file_type = ""

    try:
        if filename.endswith('.docx') or filename.endswith('.doc'):
            extracted_text = extract_from_docx(content)
            file_type = "word"
        elif filename.endswith('.xlsx') or filename.endswith('.xls'):
            extracted_text = extract_from_xlsx(content)
            file_type = "excel"
        elif filename.endswith('.pdf'):
            extracted_text = extract_from_pdf(content)
            file_type = "pdf"
        elif filename.endswith('.png') or filename.endswith('.jpg') or filename.endswith('.jpeg'):
            extracted_text = await extract_from_image(content, filename)
            file_type = "image"
        else:
            extracted_text = f"[Format file tidak didukung: {filename}]"
            file_type = "unknown"
    except Exception as e:
        extracted_text = f"[ERROR membaca file: {str(e)}]"
        file_type = "error"

    extracted_text = truncate_extracted_text(extracted_text)

    has_error = extracted_text.startswith("[ERROR") or extracted_text.startswith("[Gagal") or extracted_text.startswith("[Rate limit")
    char_count = len(extracted_text)

    return {
        "filename": file.filename,
        "file_type": file_type,
        "extracted_text": extracted_text,
        "char_count": char_count,
        "success": not has_error and char_count > 0
    }


# ==================== GROQ CALLER ====================
async def call_router(prompt, sys_prompt, hist=[], max_tokens=300, temperature=0.1):
    print(f"🤖 [GROQ CALL] Model: {TARGET_MODEL} | Tokens: {max_tokens} | Temp: {temperature}")
    if not GROQ_API_KEY or GROQ_API_KEY == "MASUKKAN_API_KEY_GROQ_ANDA_DI_SINI":
        return "ERROR_UNAUTHORIZED"

    headers = {"Authorization": f"Bearer {GROQ_API_KEY}", "Content-Type": "application/json"}
    messages = [{"role": "system", "content": sys_prompt}] + hist + [{"role": "user", "content": prompt}]

    async with httpx.AsyncClient(timeout=60.0) as client:
        try:
            res = await client.post(
                "https://api.groq.com/openai/v1/chat/completions",
                json={"model": TARGET_MODEL, "messages": messages, "max_tokens": max_tokens, "temperature": temperature, "top_p": 0.85},
                headers=headers,
            )
            if res.status_code == 200:
                res_json = res.json()
                if "choices" in res_json and len(res_json["choices"]) > 0:
                    return res_json["choices"][0]["message"]["content"]
                return "ERROR_CONNECTION"
            elif res.status_code == 429:
                print("💥 [RATE LIMIT] Hit Groq rate limit. Retrying might be needed.")
                return "ERROR_LIMIT"
            else:
                print(f"💥 [API ERROR] Status: {res.status_code} | Body: {res.text[:200]}")
                return "ERROR_CONNECTION"
        except Exception as e:
            print(f"💥 [CRITICAL ERROR]: {str(e)}")
            return "ERROR_CONNECTION"


# ==================== SIMULATION ENGINE CORE ====================

class SimulationEngine:
    def __init__(self, domain_analysis, original_problem):
        self.domain_analysis = domain_analysis
        self.original_problem = original_problem
        self.variables = self._extract_variables()
        self.entities = self._extract_entities()
        print(f"🔧 [ENGINE INIT] Variables: {len(self.variables)} → {list(self.variables.keys())}")
        print(f"🔧 [ENGINE INIT] Entities: {self.entities}")

    @staticmethod
    def _clean_var_name(raw_name):
        name = raw_name.strip()
        name = re.sub(r'^[\-\*\•\s]+', '', name)
        name = re.sub(r'^\d+[\.\)]\s*', '', name)
        name = name.strip().replace(' ', '_')
        name = name.strip('_')
        return name

    @staticmethod
    def _parse_number(raw):
        if not raw:
            return 0.0
        cleaned = raw.replace('Rp', '').replace('rp', '').replace('IDR', '').replace('Rupiah', '').replace('rupiah', '').strip()
        cleaned = re.sub(r'[^\d.,]', '', cleaned)
        if not cleaned:
            return 0.0

        if ',' in cleaned and '.' in cleaned:
            last_comma = cleaned.rfind(',')
            last_period = cleaned.rfind('.')
            if last_comma > last_period:
                integer_part = cleaned[:last_comma].replace('.', '')
                decimal_part = cleaned[last_comma + 1:]
                try: return float(f"{integer_part}.{decimal_part}")
                except ValueError: pass
            else:
                integer_part = cleaned[:last_period].replace(',', '')
                decimal_part = cleaned[last_period + 1:]
                try: return float(f"{integer_part}.{decimal_part}")
                except ValueError: pass

        if '.' in cleaned and ',' not in cleaned:
            parts = cleaned.split('.')
            if len(parts) > 1 and all(len(p) == 3 for p in parts[1:]):
                return float(''.join(parts))
            if len(parts) == 2 and len(parts[1]) == 3 and len(parts[0]) <= 3:
                return float(''.join(parts))
            try: return float(cleaned)
            except ValueError: return float(cleaned.replace('.', ''))

        if ',' in cleaned and '.' not in cleaned:
            parts = cleaned.split(',')
            if len(parts) == 2 and len(parts[1]) <= 2:
                try: return float(f"{parts[0]}.{parts[1]}")
                except ValueError: pass
            try: return float(cleaned.replace(',', ''))
            except ValueError: return 0.0

        try: return float(cleaned)
        except ValueError: return 0.0

    def _extract_entities(self):
        entities = {'name': None, 'count': None}
        ent_match = re.search(r'ENTITAS_UTAMA\s*:\s*(.+)', self.domain_analysis, re.IGNORECASE)
        if ent_match:
            entities['name'] = ent_match.group(1).strip()

        count_match = re.search(r'JUMLAH_ENTITAS\s*:\s*([\d.,]+|NULL|null)', self.domain_analysis, re.IGNORECASE)
        if count_match:
            val = count_match.group(1).strip()
            if val.upper() != 'NULL' and val:
                parsed = self._parse_number(val)
                if parsed > 0:
                    entities['count'] = parsed

        if entities['count'] is None:
            count_patterns = [
                r'(\d[\d.,]*)\s*(karyawan|worker|pegawai|orang)',
                r'(\d[\d.,]*)\s*(unit|mesin|cabang|tiket|barang|aplikasi)',
                r'(\d[\d.,]*)\s*(siswa|mahasiswa|guru|dosen)',
                r'(\d[\d.,]*)\s*(pasien|dokter|perawat)',
                r'(\d[\d.,]*)\s*(server|endpoint|device|perangkat|sistem)',
            ]
            for pattern in count_patterns:
                m = re.search(pattern, self.original_problem, re.IGNORECASE)
                if m:
                    entities['count'] = self._parse_number(m.group(1))
                    if not entities['name']:
                        entities['name'] = m.group(2).capitalize()
                    break

        if entities['count'] is None:
            if entities['name']:
                entities['count'] = 1
            else:
                domain_match = re.search(r'DOMAIN\s*:\s*(.+)', self.domain_analysis, re.IGNORECASE)
                if domain_match:
                    entities['name'] = domain_match.group(1).strip()
                else:
                    entities['name'] = "Unit Operasional"
                entities['count'] = 1
            print(f"🔧 [ENTITY FALLBACK] Inferred: {entities['name']}, count=1")

        return entities

    def _extract_variables(self):
        variables = {}
        var_section = re.search(r'VARIABEL_KUANTITATIF\s*:\s*(.*?)(?=HUBUNGAN_MATEMATIS|TARGET_USER|CONSTRAINT|METRIK_RELEVAN|\Z)', self.domain_analysis, re.DOTALL | re.IGNORECASE)

        if var_section:
            content = var_section.group(1).strip()
            for line in content.split('\n'):
                line_stripped = line.strip()
                if not line_stripped or line_stripped.startswith('#') or line_stripped.startswith('==='):
                    continue

                line_clean = re.sub(r'^[\-\*\•]\s*', '', line_stripped).strip()
                line_clean = re.sub(r'^\d+[\.\)]\s*', '', line_clean).strip()

                if not line_clean or '=' not in line_clean:
                    continue

                null_patterns = [r'=\s*NULL', r'=\s*null', r'=\s*N/A', r'=\s*na', r'=\s*TBD', r'=\s*-\s*$', r'=\s*\?\s*$']
                is_null = any(re.search(p, line_clean, re.IGNORECASE) for p in null_patterns)
                if is_null:
                    continue

                m = re.match(r'(.+?)\s*=\s*([\d.,]+)\s*(.*?)\s*\|\s*SKALA\s*:\s*(\S+)\s*\|\s*DESKRIPSI\s*:\s*(.+)', line_clean, re.IGNORECASE)
                if m:
                    name = self._clean_var_name(m.group(1))
                    value = self._parse_number(m.group(2).strip())
                    unit = m.group(3).strip()
                    scale = m.group(4).strip().upper()
                    desc = m.group(5).strip()
                    if scale in ('NULL', 'N/A', '-'): scale = 'TOTAL_GLOBAL'
                    if unit.upper() in ('NULL', 'N/A', '-'): unit = ''
                    if value > 0 and name:
                        variables[name] = {'value': value, 'unit': unit, 'scale': scale, 'description': desc or name.replace('_', ' '), 'original_value': value}
                    continue

                m = re.match(r'(.+?)\s*=\s*([\d.,]+)\s*(.*)', line_clean)
                if m:
                    name = self._clean_var_name(m.group(1))
                    raw_val = m.group(2).strip()
                    rest = m.group(3).strip()
                    value = self._parse_number(raw_val)

                    if value <= 0 or not name:
                        continue

                    scale = 'TOTAL_GLOBAL'
                    desc = ''
                    unit = rest

                    skala_m = re.search(r'SKALA\s*:\s*(\S+)', rest, re.IGNORECASE)
                    if skala_m:
                        scale = skala_m.group(1).strip().upper()
                        if scale in ('NULL', 'N/A', '-'): scale = 'TOTAL_GLOBAL'
                    desc_m = re.search(r'DESKRIPSI\s*:\s*(.+)', rest, re.IGNORECASE)
                    if desc_m: desc = desc_m.group(1).strip()

                    unit = re.sub(r'\|?\s*SKALA\s*:.*', '', unit, flags=re.IGNORECASE).strip()
                    unit = re.sub(r'\|?\s*DESKRIPSI\s*:.*', '', unit, flags=re.IGNORECASE).strip()
                    unit = unit.rstrip('|').strip()
                    if unit.upper() in ('NULL', 'N/A', '-'): unit = ''

                    name_lower = name.lower()
                    if any(kw in name_lower for kw in ['per_karyawan', 'per_orang', 'per_unit', 'per_mesin', 'per_cabang']): scale = 'PER_ENTITAS'
                    if any(kw in name_lower for kw in ['per_hari', 'per_jam', 'per_bulan', 'harian']): scale = 'RATE_WAKTU'
                    if '%' in unit or 'persen' in name_lower or 'penurunan' in name_lower:
                        if not unit: unit = '%'

                    variables[name] = {'value': value, 'unit': unit, 'scale': scale, 'description': desc or name.replace('_', ' '), 'original_value': value}
                    continue

        rel_section = re.search(r'HUBUNGAN_MATEMATIS\s*:\s*(.*?)(?=TARGET_USER|CONSTRAINT|METRIK_RELEVAN|\Z)', self.domain_analysis, re.DOTALL | re.IGNORECASE)
        if rel_section:
            for line in rel_section.group(1).strip().split('\n'):
                line_clean = re.sub(r'^[\-\*\•]\s*', '', line.strip()).strip()
                if not line_clean or '=' not in line_clean or 'NULL' in line_clean.upper(): continue
                m = re.match(r'(\w[\w_]*)\s*=\s*(.+?)=\s*([\d.,]+)', line_clean)
                if m:
                    total_name = self._clean_var_name(m.group(1))
                    if total_name not in variables:
                        total_val = self._parse_number(m.group(3).strip())
                        if total_val > 0:
                            variables[total_name] = {'value': total_val, 'unit': '', 'scale': 'TOTAL_GLOBAL', 'description': f"Derived: {m.group(2).strip()}", 'original_value': total_val}

        if not variables:
            variables = self._extract_from_problem()
        return variables

    def _extract_from_problem(self):
        variables = {}
        text = self.original_problem
        idx = 0
        for m in re.finditer(r'(gaji|biaya|harga|upah|pendapatan|revenue|pengeluaran|modal|investasi|profit|laba|rugi|cost|salary|saldo|anggaran|dana)\s*(?:per\s*(\w+)\s*)?\s*(?:adalah|sebesar|sejumlah|=|:)?\s*(?:Rp\.?\s*)?([\d.,]+)', text, re.IGNORECASE):
            label = m.group(1).strip().lower()
            per_what = m.group(2)
            raw_val = m.group(3).strip()
            value = self._parse_number(raw_val)
            if value <= 0: continue
            scale = 'PER_ENTITAS' if per_what else 'TOTAL_GLOBAL'
            var_name = f"{label}" if not per_what else f"{label}_per_{per_what}"
            var_name = re.sub(r'\s+', '_', var_name)
            if var_name in variables: var_name = f"{var_name}_{idx}"
            variables[var_name] = {'value': value, 'unit': 'Rp', 'scale': scale, 'description': f"{label}" + (f" per {per_what}" if per_what else ""), 'original_value': value}
            idx += 1

        for m in re.finditer(r'(\d[\d.,]*)\s*(karyawan|worker|pegawai|orang|unit|pieces|mesin|cabang|tiket|barang|siswa|mahasiswa|pasien|dokter|perawat|guru|dosen|server|aplikasi|sistem)', text, re.IGNORECASE):
            raw_val = m.group(1).strip()
            unit_name = m.group(2).strip().lower()
            value = self._parse_number(raw_val)
            if value <= 0: continue
            var_name = f"jumlah_{unit_name}"
            if var_name not in variables:
                variables[var_name] = {'value': value, 'unit': unit_name, 'scale': 'TOTAL_GLOBAL', 'description': f"Jumlah {unit_name}", 'original_value': value}
                idx += 1

        for m in re.finditer(r'([\d.,]+)\s*(\w+)\s*(?:per|/)\s*(hari|day|bulan|month|tahun|year|jam|hour|minggu|week)', text, re.IGNORECASE):
            raw_val = m.group(1).strip()
            unit_qty = m.group(2).strip()
            time_unit = m.group(3).strip().lower()
            value = self._parse_number(raw_val)
            if value <= 0: continue
            var_name = f"{unit_qty}_per_{time_unit}"
            var_name = re.sub(r'\s+', '_', var_name)
            if var_name not in variables:
                variables[var_name] = {'value': value, 'unit': f"{unit_qty}/{time_unit}", 'scale': 'RATE_WAKTU', 'description': f"{unit_qty} per {time_unit}", 'original_value': value}
                idx += 1

        for m in re.finditer(r'([\d.,]+)\s*%', text):
            raw_val = m.group(1).strip()
            value = self._parse_number(raw_val)
            if 0 < value <= 100:
                var_name = f"persentase_{idx}"
                variables[var_name] = {'value': value, 'unit': '%', 'scale': 'TOTAL_GLOBAL', 'description': f"Persentase {idx+1}", 'original_value': value}
                idx += 1

        if not variables:
            for m in re.finditer(r'(?:Rp\.?\s*)?([\d]{4,}[\d.,]*)', text):
                raw_val = m.group(1).strip()
                value = self._parse_number(raw_val)
                if value >= 1000:
                    var_name = f"nilai_{idx}"
                    variables[var_name] = {'value': value, 'unit': 'Rp', 'scale': 'TOTAL_GLOBAL', 'description': f"Nilai keuangan {idx+1}", 'original_value': value}
                    idx += 1
        return variables

    def get_simulatable_variables(self):
        simulatable = []
        for name, data in self.variables.items():
            value = data["value"]
            if not isinstance(value, (int, float)) or value <= 0: continue
            if name.startswith("total_") or name.startswith("jumlah_"): continue
            simulatable.append({
                "name": name, "current_value": value, "unit": data["unit"],
                "scale": data["scale"], "description": data["description"],
                "min": round(value * 0.5, 2), "max": round(value * 2.0, 2),
                "step": self._calculate_step(value)
            })
        if not simulatable:
            for name, data in self.variables.items():
                value = data["value"]
                if isinstance(value, (int, float)) and value > 0:
                    simulatable.append({
                        "name": name, "current_value": value, "unit": data["unit"],
                        "scale": data["scale"], "description": data["description"],
                        "min": round(value * 0.5, 2), "max": round(value * 2.0, 2),
                        "step": self._calculate_step(value)
                    })
        return simulatable

    def _calculate_step(self, value):
        if value >= 1000000: return 50000
        elif value >= 100000: return 10000
        elif value >= 10000: return 1000
        elif value >= 1000: return 100
        elif value >= 100: return 10
        elif value >= 10: return 1
        elif value >= 1: return 0.1
        else: return 0.01

    def run_scenario(self, adjusted_params):
        working_vars = {k: v['value'] for k, v in self.variables.items()}
        working_vars.update(adjusted_params)
        entity_count = self.entities.get('count') or 1

        results = {
            'scenario_params': adjusted_params,
            'base_calculations': {},
            'scenario_comparison': {},
            'kpis': {},
            'operational_impact': {}
        }

        for var_name, var_data in self.variables.items():
            if var_data['scale'] in ['PER_ENTITAS', 'PER_UNIT'] and entity_count:
                total_key = f"total_{var_name}"
                val = working_vars.get(var_name, var_data['value'])
                total_val = val * entity_count
                results['base_calculations'][total_key] = {
                    'formula': f"{var_name} × JUMLAH_ENTITAS",
                    'substitution': f"{val} × {entity_count}",
                    'result': round(total_val, 2), 'unit': var_data['unit']
                }

        for var_name, var_data in self.variables.items():
            if var_data['scale'] == 'RATE_WAKTU':
                val = working_vars.get(var_name, var_data['value'])
                unit = var_data.get('unit', '').lower()
                if 'hari' in unit or 'day' in unit or 'harian' in var_name.lower():
                    results['base_calculations'][f"total_bulanan_{var_name}"] = {
                        'formula': f"{var_name} × 30 hari",
                        'substitution': f"{val} × 30",
                        'result': round(val * 30, 2),
                        'unit': var_data['unit'].replace('/hari', '/bulan').replace('/day', '/month')
                    }

        financial_vars = {}
        percentage_vars = {}
        for vname, vdata in self.variables.items():
            unit = vdata.get('unit', '').lower()
            val = working_vars.get(vname, vdata['value'])
            if '%' in unit: percentage_vars[vname] = val
            elif any(c in unit for c in ['rp', 'rupiah', '$', 'usd', 'idr']) or any(c in vname.lower() for c in ['biaya', 'cost', 'harga', 'gaji', 'saldo', 'anggaran', 'dana', 'modal', 'investasi', 'pendapatan', 'revenue']): financial_vars[vname] = val

        if len(financial_vars) >= 2:
            fin_names = list(financial_vars.keys())
            for i in range(len(fin_names)):
                for j in range(i + 1, len(fin_names)):
                    n1, n2 = fin_names[i], fin_names[j]
                    v1, v2 = financial_vars[n1], financial_vars[n2]
                    if v2 > 0:
                        results['base_calculations'][f"rasio_{n1}_terhadap_{n2}"] = {
                            'formula': f"({n1} ÷ {n2}) × 100%", 'substitution': f"({v1:,.0f} ÷ {v2:,.0f}) × 100%", 'result': round((v1 / v2) * 100, 2), 'unit': '%'
                        }

        if financial_vars and percentage_vars:
            for p_name, p_val in percentage_vars.items():
                for f_name, f_val in financial_vars.items():
                    impact_val = f_val * (p_val / 100)
                    results['base_calculations'][f"dampak_{p_name}_pada_{f_name}"] = {
                        'formula': f"{f_name} × ({p_name}/100)", 'substitution': f"{f_val:,.0f} × ({p_val}/100)", 'result': round(impact_val, 2), 'unit': self.variables[f_name].get('unit', '')
                    }

        budget_keys = [k for k in financial_vars if any(b in k.lower() for b in ["anggaran", "dana", "budget", "biaya", "modal"])]
        loss_keys = [k for k in percentage_vars if any(l in k.lower() for l in ["penurunan", "loss", "rugi", "decrease"])]
        fund_keys = [k for k in financial_vars if any(f in k.lower() for f in ["saldo", "revenue", "pendapatan", "omzet"])]

        if budget_keys and loss_keys and fund_keys:
            budget_name, loss_name, base_name = budget_keys[0], loss_keys[0], fund_keys[0]
            budget_val, loss_pct, base_val = financial_vars[budget_name], percentage_vars[loss_name], financial_vars[base_name]
            loss_amount = base_val * (loss_pct / 100)
            results['base_calculations'][f"sisa_setelah_{loss_name}"] = {
                'formula': f"{base_name} - ({base_name} × {loss_name}/100)", 'substitution': f"{base_val:,.0f} - ({base_val:,.0f} × {loss_pct}/100)", 'result': round(base_val - loss_amount, 2), 'unit': self.variables[base_name].get('unit', '')
            }
            results['base_calculations'][f"kecukupan_{budget_name}"] = {
                'formula': f"({budget_name} ÷ kerugian) × 100%", 'substitution': f"({budget_val:,.0f} ÷ {loss_amount:,.0f}) × 100%", 'result': round((budget_val / loss_amount) * 100, 2) if loss_amount > 0 else 0, 'unit': '%'
            }

        for param_name, new_value in adjusted_params.items():
            if param_name in self.variables:
                old_value = self.variables[param_name]['original_value']
                pct_change = ((new_value - old_value) / old_value * 100) if old_value != 0 else 0
                results['scenario_comparison'][param_name] = {
                    'original': old_value, 'adjusted': new_value, 'change_percent': round(pct_change, 2), 'impact_direction': 'increase' if pct_change > 0 else 'decrease'
                }
                if self.variables[param_name]['scale'] in ['PER_ENTITAS', 'PER_UNIT'] and entity_count:
                    old_total, new_total = old_value * entity_count, new_value * entity_count
                    total_pct = ((new_total - old_total) / old_total * 100) if old_total != 0 else 0
                    results['scenario_comparison'][f"total_{param_name}"] = {
                        'original': old_total, 'adjusted': new_total, 'change_percent': round(total_pct, 2), 'impact_direction': 'increase' if total_pct > 0 else 'decrease'
                    }

        results['kpis'] = self._derive_kpis(working_vars, entity_count, results)
        results['operational_impact'] = self._assess_operational_impact(results)
        return results

    def _derive_kpis(self, vars_dict, entity_count, results):
        kpis = {}
        financial_vars, percentage_vars = {}, {}
        for vname, vdata in self.variables.items():
            unit, val = vdata.get('unit', '').lower(), vars_dict.get(vname, vdata['value'])
            if '%' in unit: percentage_vars[vname] = val
            elif any(c in unit for c in ['rp', 'rupiah', '$', 'usd', 'idr']) or any(c in vname.lower() for c in ['biaya', 'cost', 'harga', 'gaji', 'saldo', 'anggaran', 'dana', 'modal', 'investasi', 'pendapatan', 'revenue']): financial_vars[vname] = val

        budget_keys = [k for k in financial_vars if any(b in k.lower() for b in ["anggaran", "dana", "budget", "biaya", "modal"])]
        fund_keys = [k for k in financial_vars if any(f in k.lower() for f in ["saldo", "revenue", "pendapatan", "omzet"])]

        if budget_keys and fund_keys:
            budget_val, fund_val = financial_vars[budget_keys[0]], financial_vars[fund_keys[0]]
            if fund_val > 0:
                kpis["budget_coverage_ratio"] = {'value': round((budget_val / fund_val) * 100, 2), 'formula': f"({budget_keys[0]} ÷ {fund_keys[0]}) × 100%", 'unit': '%'}

        loss_pct_keys = [k for k in percentage_vars if any(l in k.lower() for l in ["penurunan", "loss", "rugi", "decrease", "kerugian"])]
        if fund_keys and loss_pct_keys:
            fund_val, loss_pct = financial_vars[fund_keys[0]], percentage_vars[loss_pct_keys[0]]
            loss_amount = fund_val * (loss_pct / 100)
            kpis["estimasi_kerugian"] = {'value': round(loss_amount, 2), 'formula': f"{fund_keys[0]} × ({loss_pct_keys[0]}/100)", 'unit': 'Rp'}
            if budget_keys:
                budget_val = financial_vars[budget_keys[0]]
                if loss_amount > 0:
                    kpis["kecukupan_anggaran"] = {'value': round((budget_val / loss_amount) * 100, 2), 'formula': f"({budget_keys[0]} ÷ estimasi_kerugian) × 100%", 'unit': '%'}
                    kpis["surplus_defisit_anggaran"] = {'value': round(budget_val - loss_amount, 2), 'formula': f"{budget_keys[0]} - estimasi_kerugian", 'unit': 'Rp'}

        if fund_keys and loss_pct_keys:
            fund_val, loss_pct = financial_vars[fund_keys[0]], percentage_vars[loss_pct_keys[0]]
            kpis["sisa_kapasitas"] = {'value': round(fund_val * (1 - loss_pct / 100), 2), 'formula': f"{fund_keys[0]} × (1 - {loss_pct_keys[0]}/100)", 'unit': 'Rp'}

        risk_pct_keys = [k for k in percentage_vars if any(r in k.lower() for r in ["risiko", "risk", "penurunan", "loss", "kerugian", "kegagalan"])]
        if risk_pct_keys:
            risk_val = percentage_vars[risk_pct_keys[0]]
            risk_level = "Rendah" if risk_val < 20 else ("Sedang" if risk_val < 50 else "Tinggi")
            kpis["risk_score"] = {'value': risk_val, 'formula': f"{risk_pct_keys[0]}", 'unit': f"% ({risk_level})"}

        return kpis

    def _assess_operational_impact(self, results):
        impact = {'risk_level': 'normal', 'warnings': [], 'recommendations': []}
        for param, comp in results['scenario_comparison'].items():
            if isinstance(comp, dict) and 'change_percent' in comp:
                pct = abs(comp.get('change_percent', 0))
                if pct > 100:
                    impact['warnings'].append(f"Perubahan {param} sebesar {comp['change_percent']}% terlalu ekstrem")
                    impact['risk_level'] = 'high'
                elif pct > 50:
                    impact['risk_level'] = 'elevated'
                    impact['warnings'].append(f"Perubahan {param} sebesar {comp['change_percent']}% cukup signifikan")

        if results.get('kpis', {}).get('kecukupan_anggaran'):
            sufficiency = results['kpis']['kecukupan_anggaran']['value']
            if sufficiency < 100:
                impact['risk_level'] = 'high'
                impact['warnings'].append(f"Anggaran tidak mencukupi kerugian (coverage: {sufficiency:.1f}%)")
                impact['recommendations'].append("Pertimbangkan menambah anggaran darurat atau mengurangi scope")
            elif sufficiency < 150:
                impact['risk_level'] = 'elevated'
                impact['warnings'].append(f"Margin anggaran tipis (coverage: {sufficiency:.1f}%)")

        if impact['risk_level'] == 'normal' and not impact['warnings']:
            impact['recommendations'].append("Skenario dalam batas wajar, siap untuk implementasi")
        return impact


# ==================== WEBSOCKET HANDLER ====================

@app.websocket("/ws")
async def websocket_endpoint(websocket: WebSocket):
    await websocket.accept()
    print("\n⚡ [WEBSOCKET] Connected.")
    history = []
    original_problem = ""
    domain_analysis_cache = ""
    simulation_engine = None
    current_phase = "idle"

    try:
        while True:
            data = json.loads(await websocket.receive_text())

            if data["type"] == "start":
                problem = data["problem"]
                original_problem = problem
                history.clear()
                current_phase = "deliberation"
                await websocket.send_json({"step": "roles", "progress": 5})

                domain_analysis = await call_router(
                    prompt=f"""Analisis permasalahan berikut dan ekstrak struktur kuantitatifnya.

PERMASALAHAN:
{problem}

FORMAT OUTPUT (WAJIB IKUTI PERSIS):
DOMAIN: [nama domain]
ENTITAS_UTAMA: [nama objek utama yang dianalisis]
JUMLAH_ENTITAS: [angka total atau NULL]

VARIABEL_KUANTITATIF:
- nama_variabel=3000000 Rp | SKALA:PER_ENTITAS | DESKRIPSI:Gaji per karyawan
- nama_variabel_lain=50000000 Rp | SKALA:TOTAL_GLOBAL | DESKRIPSI:Biaya operasional total

HUBUNGAN_MATEMATIS:
- total_gaji = gaji_per_karyawan × JUMLAH_ENTITAS = 30000000

TARGET_USER: [tujuan user]
CONSTRAINT: [keterbatasan]
METRIK_RELEVAN: [metrik evaluasi]

ATURAN PENTING:
1. Tulis angka TANPA titik pemisah ribuan (3000000 bukan 3.000.000)
2. Gunakan nama_variabel dengan underscore tanpa spasi
3. Ekstrak HANYA angka yang disebutkan di permasalahan
4. JANGAN buat variabel dengan nilai NULL — jika angka tidak ada, JANGAN tulis variabelnya
5. SKALA wajib salah satu: PER_ENTITAS, TOTAL_GLOBAL, RATE_WAKTU
6. Persentase (seperti 30%, 100%) tetap ditulis dengan unit % dan SKALA:TOTAL_GLOBAL
7. Jika teks mengandung ISI DOKUMEN, baca dan ekstrak semua data kuantitatif dari dokumen tersebut""",
                    sys_prompt="You are a structured data extraction engine. NEVER invent numbers. NEVER create NULL variables. Follow the format EXACTLY. If the input contains document content, extract quantitative data from it.",
                    max_tokens=600, temperature=0.05
                )

                domain_analysis_cache = domain_analysis
                print(f"📊 [DOMAIN ANALYSIS] {domain_analysis[:500]}...")
                simulation_engine = SimulationEngine(domain_analysis, problem)

                await websocket.send_json({"step": "roles", "progress": 15})

                roles_raw = await call_router(
                    f"Based on this problem: {domain_analysis[:500]}. List exactly 4 professional expert roles (comma separated, max 5 words each). ONLY return the 4 role names separated by commas.",
                    "System", max_tokens=80, temperature=0.1
                )
                roles = parse_roles(roles_raw)
                print(f"🎭 [ROLES] {roles}")

                current_progress = 20
                for role in roles:
                    await asyncio.sleep(2.5)
                    expert_prompt = f"""KONTEKS PERMASALAHAN USER:
{original_problem}

STRUKTUR KUANTITATIF:
{domain_analysis}

TUGAS ANDA sebagai {role}:
ATURAN KRITIS:
1. Jika PER_ENTITAS dan ada JUMLAH_ENTITAS, WAJIB kalikan.
2. GUNAKAN SEMUA parameter user. JANGAN ganti dengan angka generik.
3. RUMUS: [Nama] = [Substitusi Angka] = [Hasil] [Unit]
4. Bahasa Indonesia profesional. Maksimal 4 kalimat.
5. Jika ada ISI DOKUMEN, analisis data dari dokumen tersebut secara mendalam."""

                    reply = await call_router(prompt=expert_prompt, sys_prompt=f"You are a disciplined {role}. You ALWAYS verify scaling. You NEVER hallucinate numbers.", max_tokens=400, temperature=0.05)

                    if reply.startswith("ERROR_"):
                        await websocket.send_json({"step": "final", "content": "🚨 **Diskusi Terputus:** Token limit atau koneksi error.", "files": {"pdf": "", "ppt": ""}})
                        break

                    history.append({"role": "assistant", "content": f"[{role}]: {reply}"})
                    current_progress += 10
                    await websocket.send_json({"step": "discussing", "agent": role, "text": reply, "progress": min(current_progress, 60)})

                await websocket.send_json({"step": "discussing", "agent": "🔍 Reviewer Agent", "text": "Verifying dimensional analysis...", "progress": 65})
                formatted_history = "\n".join([m["content"] for m in history])

                review_result = await call_router(
                    prompt=f"PERMASALAHAN: {original_problem}\nSTRUKTUR: {domain_analysis}\nPAKAR: {formatted_history}\n\nVERIFIKASI: [DIM_ERROR], [SCALE_ERROR], [PARAM_MISS], [OP_RISK]. Koreksi numerik. Atau VALIDATED.",
                    sys_prompt="You are a quantitative auditor. Check dimensional consistency and scaling errors. BE EXTREMELY CONCISE.",
                    max_tokens=300, temperature=0.05
                )

                history.append({"role": "assistant", "content": f"[Reviewer Agent]: {review_result}"})
                await websocket.send_json({"step": "discussing", "agent": "🔍 Reviewer Agent", "text": review_result[:300] + "..." if len(review_result) > 300 else review_result, "progress": 70})

                await websocket.send_json({"step": "moderating", "progress": 80})
                needs_clarification = any(tag in review_result for tag in ["[SCALE_ERROR]", "[PARAM_MISS]", "[DIM_ERROR]"])

                if needs_clarification:
                    issues = [l for l in review_result.split("\n") if any(t in l for t in ["[SCALE_ERROR]", "[PARAM_MISS]", "[DIM_ERROR]"])]
                    clarification = "Verifikasi kalkulasi: " + (issues[0] if issues else "Inkonsistensi perhitungan terdeteksi.")
                    await websocket.send_json({"step": "ask_user", "text": sanitize_moderator_question(clarification), "progress": 85})
                    continue

                mod_eval = await call_router(
                    prompt=f"PERMASALAHAN: {original_problem}\nSTRUKTUR: {domain_analysis}\n\nTUGAS: Periksa parameter kritis yang missing. Jika cukup: READY. Jika tidak: 1 kalimat singkat.",
                    sys_prompt="You only ask for missing critical parameters.", max_tokens=60, temperature=0.1
                )

                clean_question = sanitize_moderator_question(mod_eval)
                if clean_question != "READY" and not clean_question.startswith("ERROR_"):
                    await websocket.send_json({"step": "ask_user", "text": clean_question, "progress": 85})
                else:
                    current_phase = "simulation"
                    simulatable_vars = simulation_engine.get_simulatable_variables()
                    domain_match = re.search(r'DOMAIN\s*:\s*([^\n]+)', domain_analysis, re.IGNORECASE)
                    domain_name = domain_match.group(1).strip() if domain_match else "General"

                    await websocket.send_json({
                        "step": "simulation_ready", "progress": 90,
                        "simulation_data": {
                            "domain": domain_name,
                            "entities": simulation_engine.entities,
                            "variables": simulatable_vars,
                            "original_problem": original_problem,
                            "expert_analysis": formatted_history
                        }
                    })

            elif data["type"] == "run_simulation":
                if not simulation_engine:
                    await websocket.send_json({"step": "error", "content": "Simulation engine not initialized."})
                    continue
                adjusted_params = data.get("adjusted_params", {})

                if not adjusted_params and simulation_engine.variables:
                    adjusted_params = {k: v['value'] for k, v in simulation_engine.variables.items()}

                scenario_results = simulation_engine.run_scenario(adjusted_params)
                await websocket.send_json({"step": "simulation_result", "progress": 92, "scenario_results": scenario_results})

            elif data["type"] == "skip_simulation":
                current_phase = "action_plan"
                await finish_discussion(original_problem, history, domain_analysis_cache, websocket, "VALIDATED", simulation_data={})

            elif data["type"] == "confirm_simulation":
                current_phase = "action_plan"
                sim_summary = data.get("simulation_summary", {})
                await finish_discussion(original_problem, history, domain_analysis_cache, websocket, "VALIDATED", simulation_data=sim_summary)

            elif data["type"] == "answer":
                history.append({"role": "user", "content": data["text"]})
                problem_for_finish = original_problem if original_problem else data["text"]
                if simulation_engine and current_phase == "deliberation":
                    current_phase = "simulation"
                    simulatable_vars = simulation_engine.get_simulatable_variables()
                    domain_match = re.search(r'DOMAIN\s*:\s*([^\n]+)', domain_analysis_cache, re.IGNORECASE)
                    domain_name = domain_match.group(1).strip() if domain_match else "General"
                    formatted_history = "\n".join([m["content"] for m in history])
                    await websocket.send_json({
                        "step": "simulation_ready", "progress": 90,
                        "simulation_data": {"domain": domain_name, "entities": simulation_engine.entities, "variables": simulatable_vars, "original_problem": original_problem, "expert_analysis": formatted_history}
                    })
                else:
                    await finish_discussion(problem_for_finish, history, "User provided additional context.", websocket, "VALIDATED")

    except WebSocketDisconnect:
        print("🔌 [WEBSOCKET] Disconnected.")
    except Exception as e:
        print(f"💥 [WS ERROR]: {str(e)}")


# ==================== FINAL COMPILATION ====================

async def finish_discussion(original_problem, history, domain_analysis, websocket, review_result="VALIDATED", simulation_data=None):
    print("✍️ [FINAL] Compiling blueprint...")
    formatted_history = "\n".join([m["content"] for m in history])

    sim_context = ""
    if simulation_data:
        readable_lines = []

        if simulation_data.get("scenario_params"):
            readable_lines.append("PARAMETER YANG DIUBAH:")
            for k, v in simulation_data["scenario_params"].items():
                readable_lines.append(f"  - {k}: {v:,.2f}" if isinstance(v, (int, float)) else f"  - {k}: {v}")

        if simulation_data.get("base_calculations"):
            readable_lines.append("\nPERHITUNGAN DASAR:")
            for k, v in simulation_data["base_calculations"].items():
                readable_lines.append(f"  - {k}: {v.get('formula', '')} = {v.get('substitution', '')} = {v.get('result', 0):,.2f} {v.get('unit', '')}")

        if simulation_data.get("scenario_comparison"):
            readable_lines.append("\nPERBANDINGAN SKENARIO:")
            for k, v in simulation_data["scenario_comparison"].items():
                if isinstance(v, dict) and "change_percent" in v:
                    direction = "naik" if v.get("impact_direction") == "increase" else "turun"
                    readable_lines.append(f"  - {k}: {v.get('original', 0):,.2f} → {v.get('adjusted', 0):,.2f} ({direction} {v.get('change_percent', 0):.1f}%)")

        if simulation_data.get("kpis"):
            readable_lines.append("\nKEY PERFORMANCE INDICATORS:")
            for k, v in simulation_data["kpis"].items():
                if isinstance(v, dict):
                    readable_lines.append(f"  - {k}: {v.get('value', 0):,.2f} {v.get('unit', '')} [Rumus: {v.get('formula', '')}]")

        if simulation_data.get("operational_impact"):
            oi = simulation_data["operational_impact"]
            readable_lines.append(f"\nTINGKAT RISIKO: {oi.get('risk_level', 'normal').upper()}")
            if oi.get("warnings"):
                readable_lines.append("PERINGATAN:")
                for w in oi["warnings"]: readable_lines.append(f"  ⚠️ {w}")
            if oi.get("recommendations"):
                readable_lines.append("REKOMENDASI:")
                for r in oi["recommendations"]: readable_lines.append(f"  💡 {r}")

        sim_context = f"""
HASIL SIMULASI YANG TELAH DIVALIDASI:
{chr(10).join(readable_lines)}

PENTING: Gunakan data di atas sebagai basis numerik. Jangan tampilkan sebagai JSON mentah atau kode program. Nyatakan dalam bentuk poin-poin atau paragraf profesional."""

    is_document_based = "ISI DOKUMEN" in original_problem or "DOKUMEN UPLOAD" in original_problem
    document_instruction = ""
    if is_document_based:
        document_instruction = """
8. Karena analisis ini berdasarkan DOKUMEN yang diupload user, PASTIKAN semua referensi data merujuk pada isi dokumen tersebut.
9. Sebutkan secara eksplisit data apa saja yang diambil dari dokumen.
"""

    final_prompt = f"""PERMASALAHAN ASLI:
{original_problem}

STRUKTUR KUANTITATIF:
{domain_analysis}

ANALISIS PAKAR:
{formatted_history}

{sim_context}

TUGAS: Susun "Strategic Action Resolution Blueprint" yang dinamis dan kuantitatif.

ATURAN WAJIB:
1. GUNAKAN HANYA angka dari permasalahan user.
2. Setiap perhitungan WAJIB format: RUMUS = [Substitusi] = [Hasil] [Unit]
3. PARAMETER LOCK: Gunakan SEMUA variabel user, jangan ganti dengan angka generik.
4. Jika ada hasil simulasi, nyatakan dalam bentuk NARASI DAN POIN-POIN PROFESIONAL.
5. DILARANG KERAS menampilkan format JSON, kode program, atau tanda kurung kurawal {{ }}.
6. TABEL PERBANDINGAN jika ada opsi.
7. REKOMENDASI FINAL berbasis numerik.
{document_instruction}
FORMAT BAB (Gunakan Markdown Heading ###):
### 1. [Analisis Situasi Domain]
Tulis dalam paragraf singkat.

### 2. [Hasil Simulasi & Skenario Terpilih]
Jelaskan menggunakan poin-poin (• atau -).

### 3. [Action Plan & SOP Implementasi]
Jelaskan langkah-langkah menggunakan poin-poin bernomor.

### 4. [Rekomendasi Strategis]
Tulis rekomendasi dalam poin-poin.

Bahasa Indonesia profesional. Tanpa salam. Tanpa format coding."""

    res = await call_router(
        prompt=final_prompt,
        sys_prompt="You are a Chief Strategic Officer. You write professional prose, bullet points, and paragraphs. You NEVER output raw JSON, code blocks, or curly braces. You format simulation results as natural language.",
        max_tokens=1800, temperature=0.05
    )

    if res.startswith("ERROR_"):
        await websocket.send_json({"step": "final", "content": "🚨 Gagal menyusun kesimpulan. Silakan coba lagi.", "files": {"pdf": "", "ppt": ""}})
        return

    generic_indicators = ["200.000", "200000", "200 pieces", "Rp 3.000.000", "10 jam", "25 pieces"]
    user_text_lower = original_problem.lower()
    suspicious = any(i in res and i not in user_text_lower for i in generic_indicators)

    if suspicious:
        res = await call_router(
            prompt=final_prompt + "\n\nPERINGATAN: HANYA gunakan angka dari permasalahan user. DILARANG MENGGUNAKAN FORMAT JSON ATAU KODE.",
            sys_prompt="Strict data integrity officer. Only user numbers. Format as prose.",
            max_tokens=1800, temperature=0.0
        )

    output_files = generate_files(res)
    await websocket.send_json({
        "step": "final", "content": res,
        "files": {"pdf": output_files["pdf"], "ppt": output_files["ppt"]},
        "slides_preview": output_files["slides_preview"],
        "progress": 100
    })
    print("🏁 [SUCCESS] Blueprint sent.\n")


# ==================== HELPER FUNCTIONS ====================

def sanitize_moderator_question(raw_text):
    raw = raw_text.strip()
    if "READY" in raw.upper(): return "READY"
    sentences = re.split(r'(?<=[.!?])\s+', raw)
    question = next((s.strip() for s in sentences if "?" in s), raw)
    if len(question) > 120: question = question[:120].rsplit(' ', 1)[0] + "?"
    return question

def parse_roles(raw_text):
    cleaned = raw_text.strip()
    roles = [r.strip() for r in cleaned.split(",") if r.strip()]
    if len(roles) < 4 or any(len(r) > 50 for r in roles):
        roles = [r.strip() for r in cleaned.split("\n") if r.strip()]
        roles = [r for r in roles if any(k in r.lower() for k in ["analyst", "manager", "engineer", "consultant", "strategist", "specialist", "expert", "advisor", "security", "officer"])]
    final_roles = []
    for r in roles[:6]:
        clean_r = re.sub(r'^[\d\.\-\*•\s]+', '', r).strip()
        clean_r = re.sub(r'\s*[\(\-].*$', '', clean_r).strip()
        if clean_r and 3 < len(clean_r) < 40: final_roles.append(clean_r)
    if len(final_roles) < 4: final_roles = ["Domain Analyst", "Financial Analyst", "Risk Analyst", "Strategy Analyst"]
    return final_roles[:4]

def set_text_style(paragraph, text, font_size=14, bold=False, color=THEME["text"], font_name="Calibri"):
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.space_after = Pt(8)
    return run

def add_shape_background(slide, shape_type, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def generate_files(content):
    print("📂 [FILE SYSTEM] Generating documents...")
    try:
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", 'B', 16)
        pdf.cell(0, 10, "STRATEGIC RESOLUTION & BLUEPRINT REPORT", ln=1, align='C')
        pdf.set_font("Arial", 'I', 10)
        pdf.cell(0, 10, "Argunex Multi-Agent Orchestration Framework", ln=1, align='C')
        pdf.ln(5)
        pdf.set_font("Arial", size=10)
        cleaned_text = content.encode('latin-1', 'replace').decode('latin-1')
        pdf.multi_cell(0, 6, cleaned_text)
        pdf.output("static/report.pdf")

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        sections = content.split("###")
        slides_data_json = []

        slide_layout = prs.slide_layouts[6]
        title_slide = prs.slides.add_slide(slide_layout)
        add_shape_background(title_slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height, THEME["primary"])
        add_shape_background(title_slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.15), THEME["accent"])
        circle = add_shape_background(title_slide, MSO_SHAPE.OVAL, Inches(10), Inches(4.5), Inches(4), Inches(4), THEME["secondary"])
        circle.fill.fore_color.brightness = 0.3

        tb = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(9), Inches(2))
        tf = tb.text_frame
        set_text_style(tf.paragraphs[0], "STRATEGIC RESOLUTION", 44, True, THEME["white"])
        set_text_style(tf.add_paragraph(), "& BLUEPRINT REPORT", 36, True, THEME["white"])

        sb = title_slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
        set_text_style(sb.text_frame.paragraphs[0], "Generated by Argunex AI Multi-Agent Committee", 16, False, THEME["light"])

        for section in sections:
            if not section.strip(): continue
            lines = section.strip().split("\n")
            slide_title = lines[0].strip().replace("**", "").replace("#", "")
            bullets = [l.strip().lstrip("-*• 1234567890.").strip() for l in lines[1:] if l.strip() and len(l.strip()) > 10]
            if not bullets: continue

            for chunk_idx in range(0, len(bullets), 4):
                chunk = bullets[chunk_idx:chunk_idx + 4]
                suffix = f" ({chunk_idx // 4 + 1})" if len(bullets) > 4 else ""

                cs = prs.slides.add_slide(slide_layout)
                add_shape_background(cs, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height, THEME["white"])
                add_shape_background(cs, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.1), THEME["primary"])
                add_shape_background(cs, MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), prs.slide_width, Inches(0.08), THEME["accent"])

                snb = cs.shapes.add_textbox(Inches(11.8), Inches(0.35), Inches(1.2), Inches(0.5))
                set_text_style(snb.text_frame.paragraphs[0], f"0{len(slides_data_json)+2}", 20, True, THEME["white"]).alignment = PP_ALIGN.RIGHT

                ttb = cs.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(10), Inches(0.7))
                set_text_style(ttb.text_frame.paragraphs[0], f"{slide_title}{suffix}", 24, True, THEME["white"])

                cbg = add_shape_background(cs, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.6), THEME["light"])
                cbg.fill.fore_color.brightness = 0.05

                bb = cs.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.3), Inches(5))
                btf = bb.text_frame
                btf.word_wrap = True
                for idx, item in enumerate(chunk):
                    p = btf.paragraphs[0] if idx == 0 else btf.add_paragraph()
                    set_text_style(p, f"▸  {item}", 15, False, THEME["text"])
                    p.space_before = Pt(12)
                    p.space_after = Pt(12)

                slides_data_json.append({"title": f"{slide_title}{suffix}", "bullets": chunk})

        cls = prs.slides.add_slide(slide_layout)
        add_shape_background(cls, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height, THEME["dark"])
        add_shape_background(cls, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.15), THEME["accent"])
        deco = add_shape_background(cls, MSO_SHAPE.OVAL, Inches(-1), Inches(-1), Inches(4), Inches(4), THEME["primary"])
        deco.fill.fore_color.brightness = 0.2

        tyb = cls.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(1.5))
        set_text_style(tyb.text_frame.paragraphs[0], "Thank You", 52, True, THEME["white"]).alignment = PP_ALIGN.CENTER

        tsb = cls.shapes.add_textbox(Inches(1), Inches(4.3), Inches(11), Inches(1))
        set_text_style(tsb.text_frame.paragraphs[0], "Questions & Discussion", 20, False, THEME["light"]).alignment = PP_ALIGN.CENTER

        add_shape_background(cls, MSO_SHAPE.RECTANGLE, Inches(4), Inches(5.3), Inches(5.333), Inches(0.03), THEME["accent"])

        cb = cls.shapes.add_textbox(Inches(1), Inches(5.6), Inches(11), Inches(0.8))
        set_text_style(cb.text_frame.paragraphs[0], "Argunex AI | Strategic Intelligence Engine", 12, False, THEME["light"]).alignment = PP_ALIGN.CENTER

        prs.save("static/report.pptx")
        print("✅ [FILE SYSTEM] PDF & PPTX generated.")
        return {"pdf": "/static/report.pdf", "ppt": "/static/report.pptx", "slides_preview": slides_data_json}
    except Exception as e:
        print(f"❌ [FILE SYSTEM ERROR]: {str(e)}")
        import traceback
        traceback.print_exc()
        return {"pdf": "", "ppt": "", "slides_preview": []}


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="127.0.0.1", port=8000)